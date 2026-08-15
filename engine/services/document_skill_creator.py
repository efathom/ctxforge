"""
Office Document Skill Creator.

Extracts text from office documents (PDF, DOCX, PPTX) and generates skills
via LLM synthesis. Optional dependencies (PyPDF2, python-docx, python-pptx)
are imported at extraction time with clear error messages.
"""
import json
import logging
import os
import re
from typing import Optional

from ctxforge.core.skill import Skill, SkillContent, SkillScope
from ctxforge.engine.prompts.document_skill import (
    DOCUMENT_SKILL_SYSTEM_PROMPT,
    build_document_skill_prompt,
)
from ctxforge.protocols.llm import ChatMessage, ILLMProvider

logger = logging.getLogger(__name__)

_KEBAB_RE = re.compile(r"^[a-z][a-z0-9-]*$")

_SUPPORTED_EXTENSIONS = frozenset({".pdf", ".docx", ".pptx"})

_FILE_TYPE_MAP = {
    ".pdf": "PDF Document",
    ".docx": "Word Document",
    ".pptx": "PowerPoint Presentation",
}


class DocumentSkillCreator:
    """Extract text from office documents and generate skills via LLM.

    Supports: PDF (.pdf), Word (.docx), PowerPoint (.pptx).
    Optional dependencies are imported at extraction time.
    """

    def __init__(
        self,
        llm_provider: ILLMProvider,
        max_chars: int = 50_000,
        model: Optional[str] = None,
    ):
        self._llm = llm_provider
        self._max_chars = max_chars
        self._model = model

    async def create_from_file(
        self,
        file_path: str,
        project_id: str = "default",
    ) -> Optional[Skill]:
        """Generate a skill from an office document.

        Args:
            file_path: Path to the document file.
            project_id: Project scope ID for the generated skill.

        Returns:
            A Skill object, or None if generation failed.

        Raises:
            FileNotFoundError: If the file does not exist.
            ValueError: If the file type is not supported.
            ImportError: If the required extraction library is not installed.
        """
        if not os.path.exists(file_path):
            raise FileNotFoundError(f"File not found: {file_path}")

        ext = os.path.splitext(file_path)[1].lower()
        if ext not in _SUPPORTED_EXTENSIONS:
            raise ValueError(
                f"Unsupported file type: {ext}. "
                f"Supported: {sorted(_SUPPORTED_EXTENSIONS)}"
            )

        document_text = self._extract_text(file_path)
        if not document_text.strip():
            logger.warning("No text content extracted from %s", file_path)
            return None

        file_name = os.path.basename(file_path)
        file_type = _FILE_TYPE_MAP.get(ext, "Unknown")

        user_prompt = build_document_skill_prompt(
            file_name=file_name,
            file_type=file_type,
            document_text=document_text,
        )

        messages = [
            ChatMessage(role="system", content=DOCUMENT_SKILL_SYSTEM_PROMPT),
            ChatMessage(role="user", content=user_prompt),
        ]

        response = await self._llm.chat(
            messages=messages,
            model=self._model,
            temperature=0.3,
            max_tokens=4096,
        )

        return self._parse_response(response.content, project_id)

    def _extract_text(self, file_path: str) -> str:
        """Dispatch to format-specific extractor based on extension."""
        ext = os.path.splitext(file_path)[1].lower()
        if ext == ".pdf":
            return self._extract_pdf(file_path)
        elif ext == ".docx":
            return self._extract_docx(file_path)
        elif ext == ".pptx":
            return self._extract_pptx(file_path)
        raise ValueError(f"Unsupported file type: {ext}")

    def _extract_pdf(self, file_path: str) -> str:
        """Extract text from PDF file using PyPDF2."""
        try:
            from PyPDF2 import PdfReader
        except ImportError:
            raise ImportError(
                "PyPDF2 is required for PDF extraction. "
                "Install with: pip install PyPDF2>=3.0.1"
            ) from None

        reader = PdfReader(file_path)
        text_parts = []
        total_chars = 0

        for page in reader.pages:
            page_text = page.extract_text() or ""
            if total_chars + len(page_text) > self._max_chars:
                remaining = self._max_chars - total_chars
                text_parts.append(page_text[:remaining])
                break
            text_parts.append(page_text)
            total_chars += len(page_text)

        return "\n\n".join(text_parts)

    def _extract_docx(self, file_path: str) -> str:
        """Extract text from Word document using python-docx."""
        try:
            from docx import Document
        except ImportError:
            raise ImportError(
                "python-docx is required for Word extraction. "
                "Install with: pip install python-docx>=1.2.0"
            ) from None

        doc = Document(file_path)
        text_parts = []
        total_chars = 0

        for para in doc.paragraphs:
            para_text = para.text.strip()
            if not para_text:
                continue
            if total_chars + len(para_text) > self._max_chars:
                remaining = self._max_chars - total_chars
                text_parts.append(para_text[:remaining])
                break
            text_parts.append(para_text)
            total_chars += len(para_text)

        for table in doc.tables:
            if total_chars >= self._max_chars:
                break
            for row in table.rows:
                row_text = " | ".join(
                    cell.text.strip() for cell in row.cells
                )
                if total_chars + len(row_text) > self._max_chars:
                    break
                text_parts.append(row_text)
                total_chars += len(row_text)

        return "\n\n".join(text_parts)

    def _extract_pptx(self, file_path: str) -> str:
        """Extract text from PowerPoint presentation using python-pptx."""
        try:
            from pptx import Presentation
        except ImportError:
            raise ImportError(
                "python-pptx is required for PowerPoint extraction. "
                "Install with: pip install python-pptx>=1.0.2"
            ) from None

        prs = Presentation(file_path)
        text_parts = []
        total_chars = 0

        for slide_num, slide in enumerate(prs.slides, 1):
            slide_texts = [f"--- Slide {slide_num} ---"]
            for shape in slide.shapes:
                if not shape.has_text_frame:
                    continue
                for paragraph in shape.text_frame.paragraphs:
                    para_text = paragraph.text.strip()
                    if para_text:
                        slide_texts.append(para_text)

            slide_content = "\n".join(slide_texts)
            if total_chars + len(slide_content) > self._max_chars:
                remaining = self._max_chars - total_chars
                text_parts.append(slide_content[:remaining])
                break
            text_parts.append(slide_content)
            total_chars += len(slide_content)

        return "\n\n".join(text_parts)

    @classmethod
    def is_supported(cls, file_path: str) -> bool:
        """Check if a file type is supported."""
        ext = os.path.splitext(file_path)[1].lower()
        return ext in _SUPPORTED_EXTENSIONS

    @classmethod
    def get_file_type(cls, file_path: str) -> str:
        """Get human-readable file type string."""
        ext = os.path.splitext(file_path)[1].lower()
        return _FILE_TYPE_MAP.get(ext, "Unknown")

    def _parse_response(
        self, raw: str, project_id: str,
    ) -> Optional[Skill]:
        """Parse LLM JSON response into a Skill object."""
        cleaned = raw.strip()
        if cleaned.startswith("```"):
            lines = cleaned.split("\n")
            lines = [line for line in lines if not line.strip().startswith("```")]
            cleaned = "\n".join(lines)

        try:
            data = json.loads(cleaned)
        except json.JSONDecodeError as exc:
            logger.warning("Failed to parse document skill JSON: %s", exc)
            return None

        name = data.get("name", "")
        if not _KEBAB_RE.match(name):
            logger.warning("Invalid skill name from LLM: '%s'", name)
            return None

        description = data.get("description", "")[:256]
        instructions = data.get("instructions", "")
        triggers = data.get("triggers", [])
        category = data.get("category", "document")
        when_to_use = data.get("when_to_use", "")
        scripts = data.get("scripts", {})
        references = data.get("references", {})

        return Skill(
            name=name,
            description=description,
            scope=SkillScope.PROJECT,
            scope_id=project_id,
            content=instructions,
            triggers=triggers,
            category=category,
            when_to_use=when_to_use,
            structured_content=SkillContent(
                instructions=instructions,
                scripts=scripts if isinstance(scripts, dict) else {},
                references=references if isinstance(references, dict) else {},
            ),
        )
